"""
Parser Module
Document parsing for various formats.
"""

from typing import Dict, Any, Union
import structlog
from abc import ABC, abstractmethod

from ..models import DocumentType

logger = structlog.get_logger(__name__)


class BaseParser(ABC):
    """Base document parser"""

    @abstractmethod
    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        """Parse document content to text"""
        pass


class TextParser(BaseParser):
    """Plain text parser"""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        if isinstance(content, bytes):
            return content.decode("utf-8", errors="ignore")
        return content


class HTMLParser(BaseParser):
    """HTML to text parser"""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        import re

        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")

        # Remove scripts and styles
        content = re.sub(r"<script[^>]*>.*?</script>", "", content, flags=re.DOTALL | re.IGNORECASE)
        content = re.sub(r"<style[^>]*>.*?</style>", "", content, flags=re.DOTALL | re.IGNORECASE)

        # Convert common elements
        content = re.sub(r"<br\s*/?>", "\n", content, flags=re.IGNORECASE)
        content = re.sub(r"</p>", "\n\n", content, flags=re.IGNORECASE)
        content = re.sub(r"</div>", "\n", content, flags=re.IGNORECASE)
        content = re.sub(r"</li>", "\n", content, flags=re.IGNORECASE)
        content = re.sub(r"<h[1-6][^>]*>", "\n\n## ", content, flags=re.IGNORECASE)
        content = re.sub(r"</h[1-6]>", "\n", content, flags=re.IGNORECASE)

        # Remove remaining tags
        content = re.sub(r"<[^>]+>", "", content)

        # Decode entities
        content = content.replace("&nbsp;", " ")
        content = content.replace("&amp;", "&")
        content = content.replace("&lt;", "<")
        content = content.replace("&gt;", ">")
        content = content.replace("&quot;", '"')
        content = content.replace("&#39;", "'")

        # Clean whitespace
        content = re.sub(r"\n\s*\n", "\n\n", content)
        content = re.sub(r" +", " ", content)

        return content.strip()


class MarkdownParser(BaseParser):
    """Markdown parser (keeps as-is for now)"""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        if isinstance(content, bytes):
            return content.decode("utf-8", errors="ignore")
        return content


class PDFParser(BaseParser):
    """PDF parser using pymupdf"""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        try:
            import fitz  # pymupdf

            if isinstance(content, str):
                content = content.encode("utf-8")

            doc = fitz.open(stream=content, filetype="pdf")
            text_parts = []

            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{text}")

            doc.close()
            return "\n\n".join(text_parts)

        except ImportError:
            logger.warning("pymupdf not installed, PDF parsing unavailable")
            return "[PDF content - requires pymupdf]"
        except Exception as e:
            logger.error("PDF parsing failed", error=str(e))
            raise


class DocxParser(BaseParser):
    """DOCX parser using python-docx"""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        try:
            from docx import Document as DocxDocument
            import io

            if isinstance(content, str):
                content = content.encode("utf-8")

            doc = DocxDocument(io.BytesIO(content))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also get text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text_parts.append(" | ".join(row_text))

            return "\n\n".join(text_parts)

        except ImportError:
            logger.warning("python-docx not installed")
            return "[DOCX content - requires python-docx]"
        except Exception as e:
            logger.error("DOCX parsing failed", error=str(e))
            raise


class CSVParser(BaseParser):
    """CSV parser — flattens rows to 'header: value' lines so each row is a
    self-describing chunk the retriever can match."""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        import csv
        import io

        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
        if not rows:
            return ""
        header = rows[0]
        out = []
        for row in rows[1:]:
            pairs = [f"{header[i] if i < len(header) else f'col{i}'}: {val}" for i, val in enumerate(row)]
            out.append(" | ".join(pairs))
        return ("Columns: " + ", ".join(header) + "\n\n" + "\n".join(out)).strip()


class JSONParser(BaseParser):
    """JSON parser — flattens nested objects/arrays to dotted 'path: value' lines."""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        import json

        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        try:
            data = json.loads(content)
        except Exception as e:
            logger.warning("JSON parse failed, treating as text", error=str(e))
            return content

        lines: list = []

        def _walk(node, path: str):
            if isinstance(node, dict):
                for k, v in node.items():
                    _walk(v, f"{path}.{k}" if path else str(k))
            elif isinstance(node, list):
                for idx, v in enumerate(node):
                    _walk(v, f"{path}[{idx}]")
            else:
                lines.append(f"{path}: {node}")

        _walk(data, "")
        return "\n".join(lines)


class XLSXParser(BaseParser):
    """XLSX parser using openpyxl — each sheet's rows become 'header: value' lines."""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        try:
            import openpyxl
            import io

            if isinstance(content, str):
                content = content.encode("utf-8")
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts: list = []
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                parts.append(f"[Sheet: {ws.title}]")
                header = [str(c) if c is not None else f"col{i}" for i, c in enumerate(rows[0])]
                for row in rows[1:]:
                    cells = [
                        f"{header[i] if i < len(header) else f'col{i}'}: {val}"
                        for i, val in enumerate(row)
                        if val is not None
                    ]
                    if cells:
                        parts.append(" | ".join(cells))
            wb.close()
            return "\n".join(parts)
        except ImportError:
            logger.warning("openpyxl not installed, XLSX parsing unavailable")
            return "[XLSX content - requires openpyxl]"
        except Exception as e:
            logger.error("XLSX parsing failed", error=str(e))
            raise


class PPTXParser(BaseParser):
    """PPTX parser using python-pptx — extracts text from every slide shape + notes."""

    async def parse(self, content: Union[bytes, str], metadata: Dict[str, Any] = None) -> str:
        try:
            from pptx import Presentation
            import io

            if isinstance(content, str):
                content = content.encode("utf-8")
            prs = Presentation(io.BytesIO(content))
            parts: list = []
            for idx, slide in enumerate(prs.slides):
                slide_parts: list = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs).strip()
                            if text:
                                slide_parts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text for c in row.cells]
                            slide_parts.append(" | ".join(cells))
                if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                    notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if notes:
                        slide_parts.append(f"Notes: {notes}")
                if slide_parts:
                    parts.append(f"[Slide {idx + 1}]\n" + "\n".join(slide_parts))
            return "\n\n".join(parts)
        except ImportError:
            logger.warning("python-pptx not installed, PPTX parsing unavailable")
            return "[PPTX content - requires python-pptx]"
        except Exception as e:
            logger.error("PPTX parsing failed", error=str(e))
            raise


# Parser registry
PARSERS: Dict[DocumentType, BaseParser] = {
    DocumentType.TEXT: TextParser(),
    DocumentType.HTML: HTMLParser(),
    DocumentType.MARKDOWN: MarkdownParser(),
    DocumentType.PDF: PDFParser(),
    DocumentType.DOCX: DocxParser(),
    DocumentType.CSV: CSVParser(),
    DocumentType.JSON: JSONParser(),
    DocumentType.XLSX: XLSXParser(),
    DocumentType.PPTX: PPTXParser(),
}

__all__ = [
    "BaseParser",
    "PARSERS",
    "TextParser",
    "HTMLParser",
    "MarkdownParser",
    "PDFParser",
    "DocxParser",
    "CSVParser",
    "JSONParser",
    "XLSXParser",
    "PPTXParser",
]
